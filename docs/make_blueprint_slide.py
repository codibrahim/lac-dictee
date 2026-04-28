"""
Generates docs/LacDictee_Technical_Blueprint.pptx — single slide.
Run: python3 docs/make_blueprint_slide.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

OUT_PATH = Path(__file__).parent / "LacDictee_Technical_Blueprint.pptx"

# ── Palette ───────────────────────────────────────────────────────────────────
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
PANEL   = RGBColor(0x1E, 0x1E, 0x38)
ACCENT  = RGBColor(0x6C, 0x63, 0xFF)
GREEN   = RGBColor(0x4C, 0xAF, 0x50)
YELLOW  = RGBColor(0xFF, 0xD7, 0x00)
ORANGE  = RGBColor(0xFF, 0x8C, 0x00)
CYAN    = RGBColor(0x00, 0xBF, 0xD8)
RED     = RGBColor(0xE5, 0x39, 0x35)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xCC, 0xCC, 0xFF)
GRAY    = RGBColor(0x88, 0x88, 0xAA)

W = Inches(13.33)
H = Inches(7.5)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color, line_color=None, radius=False):
    shape = slide.shapes.add_shape(
        5 if radius else 1,   # 5 = rounded rectangle, 1 = rectangle
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    if radius:
        shape.adjustments[0] = 0.05
    return shape


def txt(slide, text, x, y, w, h, size=12, bold=False,
        color=WHITE, align=PP_ALIGN.CENTER, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def multi_txt(slide, lines, x, y, w, h, size=11, color=WHITE,
              align=PP_ALIGN.CENTER, spacing=Pt(2)):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (line, line_color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = line_color


def arrow(slide, x, y, w, h):
    shape = slide.shapes.add_shape(13, x, y, w, h)  # 13 = right arrow
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ── Layout constants ──────────────────────────────────────────────────────────
TOP     = Inches(1.35)
BOX_H   = Inches(4.9)
ARR_W   = Inches(0.35)
ARR_H   = Inches(0.3)
ARR_Y   = TOP + (BOX_H / 2) - (ARR_H / 2)

# Column x positions and widths
IN_X  = Inches(0.2);   IN_W  = Inches(1.85)
AR1_X = IN_X + IN_W;   AR1_W = ARR_W
O1_X  = AR1_X + AR1_W; O1_W  = Inches(2.8)
AR2_X = O1_X + O1_W;   AR2_W = ARR_W
O2_X  = AR2_X + AR2_W; O2_W  = Inches(2.8)
AR3_X = O2_X + O2_W;   AR3_W = ARR_W
O3_X  = AR3_X + AR3_W; O3_W  = Inches(2.8)
AR4_X = O3_X + O3_W;   AR4_W = ARR_W
OUT_X = AR4_X + AR4_W; OUT_W = Inches(1.85)


def make_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK)

    # ── Header bar ────────────────────────────────────────────────────────────
    rect(slide, 0, 0, W, Inches(1.2), ACCENT)
    txt(slide, "Technical Blueprint",
        Inches(0.3), Inches(0.12), Inches(8), Inches(0.65),
        size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(slide, "Input  →  Process  →  Output",
        Inches(0.3), Inches(0.7), Inches(8), Inches(0.4),
        size=15, color=LIGHT, align=PP_ALIGN.LEFT)
    txt(slide, "LacDictée · April 2026",
        Inches(10), Inches(0.35), Inches(3), Inches(0.4),
        size=13, color=LIGHT, align=PP_ALIGN.RIGHT)

    # ── INPUT column ──────────────────────────────────────────────────────────
    rect(slide, IN_X, TOP, IN_W, BOX_H, PANEL, ACCENT, radius=True)
    txt(slide, "INPUT", IN_X, TOP + Inches(0.12), IN_W, Inches(0.38),
        size=13, bold=True, color=ACCENT)

    # Photo box
    rect(slide, IN_X + Inches(0.12), TOP + Inches(0.6),
         IN_W - Inches(0.24), Inches(1.7),
         RGBColor(0x28, 0x28, 0x48), CYAN, radius=True)
    multi_txt(slide, [
        ("📷", CYAN, False),
        ("Student photo", WHITE, True),
        ("JPG / PNG", GRAY, False),
    ], IN_X + Inches(0.12), TOP + Inches(0.7),
       IN_W - Inches(0.24), Inches(1.5),
       size=12, align=PP_ALIGN.CENTER)

    # Reference text box
    rect(slide, IN_X + Inches(0.12), TOP + Inches(2.55),
         IN_W - Inches(0.24), Inches(1.9),
         RGBColor(0x28, 0x28, 0x48), GREEN, radius=True)
    multi_txt(slide, [
        ("📄", GREEN, False),
        ("Reference text", WHITE, True),
        ("Type / paste", GRAY, False),
        ("PDF or TXT", GRAY, False),
    ], IN_X + Inches(0.12), TOP + Inches(2.65),
       IN_W - Inches(0.24), Inches(1.7),
       size=12, align=PP_ALIGN.CENTER)

    # ── Arrow 1 ───────────────────────────────────────────────────────────────
    arrow(slide, AR1_X, ARR_Y, AR1_W, ARR_H)

    # ── OCR box ───────────────────────────────────────────────────────────────
    rect(slide, O1_X, TOP, O1_W, BOX_H, PANEL, CYAN, radius=True)
    txt(slide, "STEP 1 — OCR",
        O1_X, TOP + Inches(0.12), O1_W, Inches(0.38),
        size=12, bold=True, color=CYAN)
    txt(slide, "Handwriting → Text",
        O1_X, TOP + Inches(0.45), O1_W, Inches(0.3),
        size=10, color=GRAY)

    ocr_engines = [
        ("★ Claude Vision",   "claude-haiku-4-5",        CYAN,   RGBColor(0x22, 0x22, 0x42)),
        ("★ Infinity-Parser", "HuggingFace · free",       YELLOW, RGBColor(0x22, 0x22, 0x42)),
        ("★ Tesseract",       "local · last resort",      GRAY,   RGBColor(0x22, 0x22, 0x42)),
    ]
    labels = ["primary", "fallback", "last resort"]
    for i, (name, sub, col, bg_col) in enumerate(ocr_engines):
        ey = TOP + Inches(0.85) + i * Inches(1.25)
        rect(slide, O1_X + Inches(0.12), ey,
             O1_W - Inches(0.24), Inches(1.1), bg_col, col, radius=True)
        multi_txt(slide, [
            (name, col, True),
            (sub, GRAY, False),
            (f"[{labels[i]}]", col, False),
        ], O1_X + Inches(0.15), ey + Inches(0.06),
           O1_W - Inches(0.3), Inches(1.0),
           size=11, align=PP_ALIGN.CENTER)

    # ── Arrow 2 ───────────────────────────────────────────────────────────────
    arrow(slide, AR2_X, ARR_Y, AR2_W, ARR_H)
    txt(slide, "student_text",
        AR2_X - Inches(0.05), ARR_Y + Inches(0.33),
        AR2_W + Inches(0.1), Inches(0.28),
        size=8, color=GRAY, align=PP_ALIGN.CENTER)

    # ── AI Correction box ─────────────────────────────────────────────────────
    rect(slide, O2_X, TOP, O2_W, BOX_H, PANEL, ACCENT, radius=True)
    txt(slide, "STEP 2 — AI Correction",
        O2_X, TOP + Inches(0.12), O2_W, Inches(0.38),
        size=12, bold=True, color=ACCENT)
    txt(slide, "Compare · Detect errors",
        O2_X, TOP + Inches(0.45), O2_W, Inches(0.3),
        size=10, color=GRAY)

    rect(slide, O2_X + Inches(0.12), TOP + Inches(0.85),
         O2_W - Inches(0.24), Inches(1.6),
         RGBColor(0x22, 0x22, 0x42), ACCENT, radius=True)
    multi_txt(slide, [
        ("Claude Sonnet 4.6", ACCENT, True),
        ("claude-sonnet-4-6", GRAY, False),
        ("French teacher persona", WHITE, False),
        ("student_text vs correct_text", GRAY, False),
    ], O2_X + Inches(0.15), TOP + Inches(0.9),
       O2_W - Inches(0.3), Inches(1.5),
       size=11, align=PP_ALIGN.CENTER)

    rect(slide, O2_X + Inches(0.12), TOP + Inches(2.6),
         O2_W - Inches(0.24), Inches(2.0),
         RGBColor(0x22, 0x22, 0x42), YELLOW, radius=True)
    multi_txt(slide, [
        ("JSON output", YELLOW, True),
        ("errors[]", WHITE, False),
        ("  wrong → correct", GRAY, False),
        ("  type · explanation", GRAY, False),
        ("score  0–100", WHITE, False),
        ("feedback text", GRAY, False),
    ], O2_X + Inches(0.15), TOP + Inches(2.65),
       O2_W - Inches(0.3), Inches(1.85),
       size=11, align=PP_ALIGN.LEFT)

    # ── Arrow 3 ───────────────────────────────────────────────────────────────
    arrow(slide, AR3_X, ARR_Y, AR3_W, ARR_H)

    # ── Report Generation box ─────────────────────────────────────────────────
    rect(slide, O3_X, TOP, O3_W, BOX_H, PANEL, GREEN, radius=True)
    txt(slide, "STEP 3 — Report",
        O3_X, TOP + Inches(0.12), O3_W, Inches(0.38),
        size=12, bold=True, color=GREEN)
    txt(slide, "Generate · Format · Export",
        O3_X, TOP + Inches(0.45), O3_W, Inches(0.3),
        size=10, color=GRAY)

    report_items = [
        ("📋 Error list",     "wrong→correct · type · explanation", GREEN),
        ("🔢 Score",          "0–100, weighted severity",            YELLOW),
        ("🎨 Visual markup",  "red / green word annotation",         CYAN),
        ("📥 PDF export",     "fpdf2 · student report",             ORANGE),
    ]
    for i, (title, sub, col) in enumerate(report_items):
        ry = TOP + Inches(0.82) + i * Inches(0.98)
        rect(slide, O3_X + Inches(0.12), ry,
             O3_W - Inches(0.24), Inches(0.85),
             RGBColor(0x22, 0x22, 0x42), col, radius=True)
        multi_txt(slide, [
            (title, col, True),
            (sub,   GRAY, False),
        ], O3_X + Inches(0.15), ry + Inches(0.05),
           O3_W - Inches(0.3), Inches(0.78),
           size=10, align=PP_ALIGN.LEFT)

    # ── Arrow 4 ───────────────────────────────────────────────────────────────
    arrow(slide, AR4_X, ARR_Y, AR4_W, ARR_H)

    # ── OUTPUT column ─────────────────────────────────────────────────────────
    rect(slide, OUT_X, TOP, OUT_W, BOX_H, PANEL, GREEN, radius=True)
    txt(slide, "OUTPUT", OUT_X, TOP + Inches(0.12), OUT_W, Inches(0.38),
        size=13, bold=True, color=GREEN)

    outputs = [
        ("🖥️",  "On-screen\nreport",    "Streamlit",       ACCENT),
        ("📥",  "PDF\ndownload",        "fpdf2",            GREEN),
        ("🗄️", "History\nlog",         "SQLite",           YELLOW),
    ]
    for i, (icon, label, tech, col) in enumerate(outputs):
        oy = TOP + Inches(0.6) + i * Inches(1.4)
        rect(slide, OUT_X + Inches(0.1), oy,
             OUT_W - Inches(0.2), Inches(1.2),
             RGBColor(0x28, 0x28, 0x48), col, radius=True)
        multi_txt(slide, [
            (icon, col, False),
            (label, WHITE, True),
            (tech,  GRAY, False),
        ], OUT_X + Inches(0.12), oy + Inches(0.06),
           OUT_W - Inches(0.24), Inches(1.1),
           size=11, align=PP_ALIGN.CENTER)

    # ── Bottom stack bar ──────────────────────────────────────────────────────
    rect(slide, 0, Inches(6.35), W, Inches(0.5), RGBColor(0x11, 0x11, 0x22))
    txt(slide,
        "Stack:  Python 3.11  ·  Streamlit  ·  Claude API (Haiku + Sonnet)  ·  "
        "pytesseract  ·  PyMuPDF  ·  fpdf2  ·  SQLite",
        Inches(0.3), Inches(6.38), Inches(12.7), Inches(0.4),
        size=11, color=GRAY, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    make_slide(prs)
    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
