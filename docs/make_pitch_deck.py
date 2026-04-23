"""
Generates docs/LacDictee_MVP_Pitch_Deck.pptx
Run: python3 docs/make_pitch_deck.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SWOT_IMG = Path(__file__).parent / "research" / "SWOT_analysis.png"
OUT_PATH  = Path(__file__).parent / "LacDictee_MVP_Pitch_Deck.pptx"

# ── Palette ───────────────────────────────────────────────────────────────────
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
ACCENT  = RGBColor(0x6C, 0x63, 0xFF)   # purple-blue
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF0, 0xF0, 0xFF)
YELLOW  = RGBColor(0xFF, 0xD7, 0x00)
GREEN   = RGBColor(0x4C, 0xAF, 0x50)
RED     = RGBColor(0xE5, 0x39, 0x35)
GRAY    = RGBColor(0xCC, 0xCC, 0xCC)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=24, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def header_bar(slide, title: str, subtitle: str = ""):
    add_rect(slide, 0, 0, W, Inches(1.4), ACCENT)
    add_text(slide, title, Inches(0.4), Inches(0.1), Inches(12), Inches(0.8),
             size=36, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                 size=16, color=LIGHT)


def bullet_card(slide, x, y, w, h, title, bullets,
                bg=DARK, title_color=ACCENT, bullet_color=WHITE,
                title_size=20, bullet_size=16):
    add_rect(slide, x, y, w, h, bg)
    add_text(slide, title, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.4),
             size=title_size, bold=True, color=title_color)
    box = slide.shapes.add_textbox(
        x + Inches(0.15), y + Inches(0.55),
        w - Inches(0.3), h - Inches(0.65))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(bullet_size)
        p.font.color.rgb = bullet_color
        p.space_after = Pt(4)


# ── Slide 1 — Problem & Solution ──────────────────────────────────────────────

def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
    slide_bg(slide, DARK)
    header_bar(slide,
               "Problem & Solution",
               "Why LacDictee exists")

    # Left column — Problem
    bullet_card(slide,
                x=Inches(0.3), y=Inches(1.55), w=Inches(4.1), h=Inches(4.9),
                title="The Problem",
                bullets=[
                    "Teachers correct dictée by hand",
                    "1–2 hours per class, per week",
                    "Repetitive, error-prone process",
                    "No structured feedback for students",
                ],
                bg=RGBColor(0x25, 0x25, 0x3D),
                title_color=RED)

    # Center — Arrow
    add_text(slide, "→", Inches(4.55), Inches(3.5), Inches(0.7), Inches(0.7),
             size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Right column — Solution
    bullet_card(slide,
                x=Inches(5.4), y=Inches(1.55), w=Inches(4.1), h=Inches(4.9),
                title="The Solution",
                bullets=[
                    "Upload reference text + student photo",
                    "OCR reads handwriting automatically",
                    "Claude AI detects errors",
                    "Instant report with grade suggestion",
                ],
                bg=RGBColor(0x25, 0x25, 0x3D),
                title_color=GREEN)

    # Impact badge — right side
    add_rect(slide, Inches(9.75), Inches(2.2), Inches(3.2), Inches(3.1),
             RGBColor(0x6C, 0x63, 0xFF))
    add_text(slide, "⚡", Inches(9.85), Inches(2.3), Inches(3.0), Inches(0.8),
             size=36, align=PP_ALIGN.CENTER, color=WHITE)
    add_text(slide, "Up to", Inches(9.85), Inches(3.0), Inches(3.0), Inches(0.5),
             size=18, color=LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "90%", Inches(9.85), Inches(3.4), Inches(3.0), Inches(0.8),
             size=48, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, "time saved", Inches(9.85), Inches(4.1), Inches(3.0), Inches(0.5),
             size=18, color=LIGHT, align=PP_ALIGN.CENTER)

    # Bottom tagline
    add_rect(slide, 0, Inches(6.65), W, Inches(0.85), RGBColor(0x11, 0x11, 0x22))
    add_text(slide,
             "\"We turn a 2-hour correction task into a 2-minute review.\"",
             Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.7),
             size=20, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)


# ── Slide 2 — Business Model ──────────────────────────────────────────────────

def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, DARK)
    header_bar(slide, "Business Model", "Who, How, and Where we grow")

    card_y = Inches(1.6)
    card_h = Inches(2.0)
    gap    = Inches(0.18)
    cw     = Inches(3.1)

    cards = [
        ("Customers",
         ["French teachers", "Language schools", "Tutors & parents"],
         RGBColor(0x1E, 0x1E, 0x36), ACCENT),
        ("Market",
         ["EdTech & Language Learning", "Growing AI adoption", "Underserved teacher segment"],
         RGBColor(0x1E, 0x1E, 0x36), RGBColor(0x00, 0xBF, 0xD8)),
        ("Revenue",
         ["Freemium — free for individuals", "Subscription (monthly / yearly)", "School licensing (future)"],
         RGBColor(0x1E, 0x1E, 0x36), GREEN),
        ("Go-To-Market",
         ["Teacher Slack / LinkedIn groups", "Direct school outreach", "Demo-based adoption"],
         RGBColor(0x1E, 0x1E, 0x36), YELLOW),
    ]

    for i, (title, bullets, bg, tc) in enumerate(cards):
        x = Inches(0.3) + i * (cw + gap)
        bullet_card(slide, x, card_y, cw, card_h,
                    title=title, bullets=bullets,
                    bg=bg, title_color=tc, bullet_size=15)

    # Value row
    add_rect(slide, 0, Inches(3.75), W, Inches(1.55), RGBColor(0x11, 0x11, 0x22))
    add_text(slide, "Value Proposition", Inches(0.4), Inches(3.82), Inches(4), Inches(0.5),
             size=18, bold=True, color=ACCENT)

    metrics = [
        ("Before LacDictee", "1–2 hrs / week", WHITE),
        ("After LacDictee",  "2–5 minutes",    GREEN),
        ("Time Saved",       "80–90%",          YELLOW),
    ]
    for j, (label, val, col) in enumerate(metrics):
        mx = Inches(0.4) + j * Inches(4.1)
        add_text(slide, label, mx, Inches(4.0), Inches(3.8), Inches(0.35),
                 size=13, color=GRAY)
        add_text(slide, val, mx, Inches(4.3), Inches(3.8), Inches(0.7),
                 size=26, bold=True, color=col)

    # Organization
    add_rect(slide, 0, Inches(5.35), W, Inches(1.35), RGBColor(0x1A, 0x1A, 0x2E))
    add_text(slide, "Organization", Inches(0.4), Inches(5.42), Inches(3), Inches(0.4),
             size=16, bold=True, color=ACCENT)
    add_text(slide,
             "MVP:  Solo founder      |      Next:  AI engineer · Frontend dev · Education expert",
             Inches(0.4), Inches(5.82), Inches(12.5), Inches(0.6),
             size=16, color=LIGHT, align=PP_ALIGN.CENTER)


# ── Slide 3 — SWOT ────────────────────────────────────────────────────────────

def slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_bg(slide, DARK)
    header_bar(slide, "SWOT Analysis", "Honest view of where we stand")

    if SWOT_IMG.exists():
        slide.shapes.add_picture(
            str(SWOT_IMG),
            Inches(0.3), Inches(1.55),
            width=Inches(8.0))
    else:
        add_text(slide, "[SWOT image not found]",
                 Inches(0.3), Inches(2), Inches(8), Inches(1),
                 size=18, color=RED)

    # Side panel — key takeaway
    add_rect(slide, Inches(8.55), Inches(1.55), Inches(4.45), Inches(5.7),
             RGBColor(0x1E, 0x1E, 0x36))
    add_text(slide, "Key Takeaways", Inches(8.7), Inches(1.65), Inches(4.1), Inches(0.5),
             size=18, bold=True, color=ACCENT)

    takeaways = [
        ("Strength",     "Real problem, AI-powered, time-saving", GREEN),
        ("Weakness",     "OCR depends on handwriting quality",    YELLOW),
        ("Opportunity",  "EdTech growth + AI adoption wave",      RGBColor(0x00, 0xBF, 0xD8)),
        ("Threat",       "Privacy concerns + big competitors",    RED),
    ]
    for k, (label, note, col) in enumerate(takeaways):
        ty = Inches(2.3) + k * Inches(1.15)
        add_rect(slide, Inches(8.65), ty, Inches(4.15), Inches(0.95),
                 RGBColor(0x25, 0x25, 0x3D))
        add_text(slide, label, Inches(8.8), ty + Inches(0.05), Inches(3.8), Inches(0.35),
                 size=13, bold=True, color=col)
        add_text(slide, note, Inches(8.8), ty + Inches(0.37), Inches(3.8), Inches(0.5),
                 size=13, color=LIGHT)

    # Bottom bar
    add_rect(slide, 0, Inches(7.1), W, Inches(0.4), ACCENT)
    add_text(slide, "LacDictee  —  MVP Architecture Defense  —  April 24, 2026",
             Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.35),
             size=12, color=WHITE, align=PP_ALIGN.CENTER)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide1(prs)
    slide2(prs)
    slide3(prs)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")


if __name__ == "__main__":
    main()
